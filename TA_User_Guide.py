"""
สร้างสไลด์คู่มือสำหรับ TA Team — HC Request System
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ──────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x00, 0x80, 0x65)   # #008065 brand green
GREEN_LIGHT = RGBColor(0xE6, 0xF4, 0xF1)   # light green bg
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
GRAY        = RGBColor(0x64, 0x74, 0x8B)   # slate-500
GRAY_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # slate-100
INDIGO      = RGBColor(0x43, 0x38, 0xCA)   # indigo-600
TEAL        = RGBColor(0x0D, 0x94, 0x88)   # teal-600
ORANGE      = RGBColor(0xEA, 0x58, 0x0C)   # orange-600
RED         = RGBColor(0xDC, 0x26, 0x26)   # red-600
AMBER       = RGBColor(0xD9, 0x77, 0x06)   # amber-600

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout

# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=GREEN, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_pill(slide, text, x, y, color=GREEN, text_color=WHITE, size=13):
    """เพิ่ม pill/badge"""
    tw = Inches(1.6)
    th = Inches(0.36)
    r = add_rect(slide, x, y, tw, th, fill=color)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = text_color
    return r

def add_header_bar(slide, title, subtitle=None):
    """แถบหัวสีเขียว"""
    add_rect(slide, 0, 0, W, Inches(1.45), fill=GREEN)
    add_text(slide, title, Inches(0.5), Inches(0.18), Inches(11), Inches(0.7),
             size=32, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.5), Inches(0.88), Inches(10), Inches(0.45),
                 size=16, color=RGBColor(0xBB, 0xF7, 0xD0), italic=True)

def add_section_label(slide, text, x, y, color=GREEN):
    add_text(slide, text, x, y, Inches(4), Inches(0.35),
             size=11, bold=True, color=color)

def add_card(slide, x, y, w, h, fill=GRAY_LIGHT):
    r = add_rect(slide, x, y, w, h, fill=fill)
    r.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    r.line.width = Pt(1)
    return r

def bullet_box(slide, items, x, y, w, h, size=15, color=DARK, icon="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# full green bg top half
add_rect(slide, 0, 0, W, Inches(4.2), fill=GREEN)
# bottom white
add_rect(slide, 0, Inches(4.2), W, Inches(3.3), fill=WHITE)

# logo circle
circ = slide.shapes.add_shape(9, Inches(0.55), Inches(0.45), Inches(0.9), Inches(0.9))  # oval
circ.fill.solid(); circ.fill.fore_color.rgb = WHITE
circ.line.fill.background()
add_text(slide, "HC", Inches(0.62), Inches(0.52), Inches(0.75), Inches(0.6),
         size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_text(slide, "HC Request System", Inches(0.5), Inches(1.3), Inches(12), Inches(1.1),
         size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "คู่มือการใช้งานสำหรับ TA Team", Inches(0.5), Inches(2.5), Inches(12), Inches(0.7),
         size=26, color=RGBColor(0xBB, 0xF7, 0xD0), align=PP_ALIGN.CENTER, italic=True)
add_text(slide, "Talent Acquisition · People Experience", Inches(0.5), Inches(3.1), Inches(12), Inches(0.55),
         size=16, color=RGBColor(0x6E, 0xE7, 0xB7), align=PP_ALIGN.CENTER)

# decorative line
add_rect(slide, Inches(4.5), Inches(4.05), Inches(4.33), Inches(0.04), fill=GREEN)

add_text(slide, "TA ทำอะไรได้บ้างในระบบ?", Inches(0.5), Inches(4.3), Inches(12), Inches(0.6),
         size=22, bold=True, color=DARK, align=PP_ALIGN.CENTER)

features = ["รับเคสและเปลี่ยนสถานะ",  "ดู Dashboard ภาพรวม",  "จัดการ JD Files",  "ดู Audit Log"]
colors   = [GREEN, INDIGO, TEAL, AMBER]
for i, (feat, col) in enumerate(zip(features, colors)):
    bx = Inches(0.6 + i * 3.1)
    add_card(slide, bx, Inches(5.05), Inches(2.8), Inches(0.85), fill=GRAY_LIGHT)
    add_text(slide, feat, bx + Inches(0.12), Inches(5.08), Inches(2.6), Inches(0.75),
             size=15, bold=True, color=col, align=PP_ALIGN.CENTER)

add_text(slide, "hcrequest.web.app", Inches(0.5), Inches(6.55), Inches(12), Inches(0.4),
         size=13, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Overview / Menu
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "ภาพรวมระบบ", "หน้าจอหลักที่ TA จะใช้งาน")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=GRAY_LIGHT)

menus = [
    ("📊", "Dashboard",       "ภาพรวมทุกเคส\nสถิติ pipeline TA",         INDIGO),
    ("📋", "All Requests",    "ดู request ทั้งหมด\nกรอง / ค้นหาได้",      GREEN),
    ("🗂️", "My Cases",        "เฉพาะเคสที่\nรับผิดชอบอยู่",               TEAL),
    ("📁", "JD Files",        "ไฟล์ Job Description\nแนบกับแต่ละตำแหน่ง", AMBER),
    ("🔍", "Audit Log",       "ประวัติทุกการเปลี่ยนแปลง\nใช้ตรวจสอบ",    ORANGE),
]
for i, (icon, name, desc, col) in enumerate(menus):
    bx = Inches(0.38 + i * 2.55)
    add_card(slide, bx, Inches(1.7), Inches(2.4), Inches(4.85))
    # icon circle
    ic = slide.shapes.add_shape(9, bx + Inches(0.82), Inches(2.0), Inches(0.75), Inches(0.75))
    ic.fill.solid(); ic.fill.fore_color.rgb = col
    ic.line.fill.background()
    add_text(slide, icon, bx + Inches(0.82), Inches(2.02), Inches(0.75), Inches(0.7),
             size=22, align=PP_ALIGN.CENTER)
    add_text(slide, name, bx + Inches(0.08), Inches(2.85), Inches(2.25), Inches(0.5),
             size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(slide, desc, bx + Inches(0.08), Inches(3.4), Inches(2.25), Inches(0.9),
             size=13, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "💡  Sidebar ด้านซ้ายมือมีเมนูครบทุกอัน — คลิกเพื่อสลับหน้าได้เลย",
         Inches(0.5), Inches(6.85), Inches(12), Inches(0.45),
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Status Workflow
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "Status Workflow", "ลำดับสถานะที่ TA จะเปลี่ยนในแต่ละเคส")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=WHITE)

statuses = [
    ("Open",        "📂", RGBColor(0x64,0x74,0x8B), "Manager เปิดเคสใหม่\nรอ TA รับ"),
    ("Recruiting",  "🔍", GREEN,                     "กำลังหา\nผู้สมัคร"),
    ("Interviewing","🗣️", INDIGO,                    "อยู่ระหว่าง\nสัมภาษณ์"),
    ("Offering",    "📋", AMBER,                     "ส่ง Offer\n+ บันทึก CV"),
    ("Onboarding",  "🟦", TEAL,                      "รอเริ่มงาน\n+ วันที่"),
    ("Closed",      "✅", RGBColor(0x16,0xA3,0x4A), "Onboard\nสำเร็จ"),
]

box_w = Inches(1.75)
box_h = Inches(2.2)
gap   = Inches(0.32)
start_x = Inches(0.35)
y_box = Inches(1.85)

for i, (name, icon, col, desc) in enumerate(statuses):
    bx = start_x + i * (box_w + gap)
    # shadow
    s = add_rect(slide, bx + Inches(0.04), y_box + Inches(0.04), box_w, box_h,
                 fill=RGBColor(0xE2,0xE8,0xF0))
    # card
    card = add_rect(slide, bx, y_box, box_w, box_h, fill=WHITE)
    card.line.color.rgb = col; card.line.width = Pt(2)
    # top color bar
    add_rect(slide, bx, y_box, box_w, Inches(0.12), fill=col)
    # icon
    add_text(slide, icon, bx, y_box + Inches(0.2), box_w, Inches(0.65),
             size=28, align=PP_ALIGN.CENTER)
    # name
    add_text(slide, name, bx, y_box + Inches(0.88), box_w, Inches(0.45),
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    # desc
    add_text(slide, desc, bx + Inches(0.1), y_box + Inches(1.35), box_w - Inches(0.2), Inches(0.75),
             size=12, color=GRAY, align=PP_ALIGN.CENTER)
    # arrow (ไม่ใส่ตัวสุดท้าย)
    if i < len(statuses) - 1:
        ax = bx + box_w + Inches(0.05)
        add_text(slide, "→", ax, y_box + Inches(0.75), gap, Inches(0.6),
                 size=22, bold=True, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)

# Reject note
add_rect(slide, Inches(0.35), Inches(4.4), Inches(5.5), Inches(0.7), fill=RGBColor(0xFE,0xF2,0xF2))
add_text(slide, "❌  Rejected — TA สามารถ Reject เคสได้ทุกสถานะ พร้อมกรอกเหตุผล",
         Inches(0.5), Inches(4.45), Inches(5.2), Inches(0.55), size=13, color=RED)

add_rect(slide, Inches(6.2), Inches(4.4), Inches(6.5), Inches(0.7), fill=RGBColor(0xFF,0xF7,0xED))
add_text(slide, "🚫  Cancelled — เฉพาะ Admin เท่านั้น (ยกเลิกโดย Manager/HR)",
         Inches(6.35), Inches(4.45), Inches(6.1), Inches(0.55), size=13, color=ORANGE)

add_text(slide, "💡  TA เปลี่ยน status ได้จาก dropdown ในตาราง All Requests หรือ My Cases",
         Inches(0.5), Inches(5.35), Inches(12), Inches(0.4),
         size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dashboard
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "📊 Dashboard", "ภาพรวมทุกเคสในระบบ — หน้าแรกที่เห็นหลัง Login")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=GRAY_LIGHT)

# Stat Cards mockup
cards_data = [
    ("Active Requests", "24", GREEN),
    ("Recruiting",      "8",  INDIGO),
    ("Offering",        "3",  AMBER),
    ("Closed (เดือนนี้)","5", TEAL),
]
for i, (label, val, col) in enumerate(cards_data):
    bx = Inches(0.35 + i * 3.25)
    add_card(slide, bx, Inches(1.65), Inches(3.0), Inches(1.3))
    add_rect(slide, bx, Inches(1.65), Inches(0.07), Inches(1.3), fill=col)
    add_text(slide, val,   bx + Inches(0.2), Inches(1.8),  Inches(2.5), Inches(0.6),
             size=36, bold=True, color=col)
    add_text(slide, label, bx + Inches(0.2), Inches(2.4),  Inches(2.5), Inches(0.4),
             size=13, color=GRAY)

# Features list
feat_y = Inches(3.15)
add_text(slide, "สิ่งที่เห็นใน Dashboard:", Inches(0.4), feat_y, Inches(5), Inches(0.4),
         size=15, bold=True, color=DARK)
items = [
    "📈  Stat Cards — จำนวนเคสแต่ละสถานะ",
    "📅  Monthly Pipeline — กราฟ timeline แต่ละเดือน",
    "👤  TA Workload — เคสของแต่ละ TA",
    "📋  ตารางเคสทั้งหมด พร้อม filter / sort",
    "⏱️  SLA Timer — นับวันที่ใช้ต่อเคส",
]
bullet_box(slide, items, Inches(0.4), Inches(3.6), Inches(6.2), Inches(2.8),
           size=14, icon="")

# Right panel tip
add_card(slide, Inches(7.0), Inches(3.15), Inches(5.9), Inches(3.0), fill=RGBColor(0xEC,0xFD,0xF5))
add_text(slide, "💡 Tips", Inches(7.2), Inches(3.3), Inches(5.5), Inches(0.4),
         size=15, bold=True, color=GREEN)
tips = [
    "กดที่ชื่อ TA ใน Workload Panel\nเพื่อ filter ดูเฉพาะเคสของ TA คนนั้น",
    "แถบสี SLA แดง = เกินกำหนด,\nเขียว = อยู่ในเวลา",
    "กด tab สถานะด้านบนตารางเพื่อ\nกรองเฉพาะ status ที่ต้องการ",
]
bullet_box(slide, tips, Inches(7.2), Inches(3.75), Inches(5.5), Inches(2.2),
           size=13, color=RGBColor(0x06,0x5F,0x46), icon="•")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — All Requests
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "📋 All Requests", "ดู request ทั้งหมดในระบบ — เปลี่ยน status ได้จากหน้านี้")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=WHITE)

# Table mock header
add_rect(slide, Inches(0.3), Inches(1.65), Inches(12.7), Inches(0.5), fill=GREEN)
headers = ["HCID", "Position", "Department", "Status", "PIC (TA)", "Open Date", "SLA"]
col_ws  = [1.4, 2.5, 2.2, 1.6, 1.8, 1.5, 1.0]
cx = Inches(0.35)
for h, cw in zip(headers, col_ws):
    add_text(slide, h, cx, Inches(1.7), Inches(cw), Inches(0.4),
             size=12, bold=True, color=WHITE)
    cx += Inches(cw)

# Table rows mock
rows = [
    ("REQ-2026-429","Marketing Manager","Marketing","Recruiting","Noon T.","27-Apr-26","12d"),
    ("REQ-2026-428","Data Analyst","Data Team","Offering","Art S.","20-Apr-26","19d"),
    ("REQ-2026-427","Logistics Officer","Logistic","Open","—","15-Apr-26","24d"),
]
row_colors = [WHITE, GRAY_LIGHT, WHITE]
for ri, (row, rc) in enumerate(zip(rows, row_colors)):
    ry = Inches(2.2 + ri * 0.55)
    add_rect(slide, Inches(0.3), ry, Inches(12.7), Inches(0.52), fill=rc)
    cx = Inches(0.35)
    for ci, (cell, cw) in enumerate(zip(row, col_ws)):
        col = DARK
        if ci == 3:  # status
            scol = {"Recruiting": INDIGO, "Offering": AMBER, "Open": GRAY}.get(cell, DARK)
            add_rect(slide, cx + Inches(0.05), ry + Inches(0.1), Inches(cw - 0.15), Inches(0.32),
                     fill=scol)
            add_text(slide, cell, cx + Inches(0.05), ry + Inches(0.12), Inches(cw - 0.15), Inches(0.3),
                     size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        else:
            add_text(slide, cell, cx, ry + Inches(0.1), Inches(cw), Inches(0.35),
                     size=12, color=col)
        cx += Inches(cw)

# Features
feat_y = Inches(4.0)
cols2 = [
    ("🔍 ค้นหา & กรอง", [
        "พิมพ์ชื่อตำแหน่ง / แผนก / HCID",
        "กดแท็บ status ด้านบนตาราง",
        "filter เฉพาะเคสของตัวเอง",
    ]),
    ("⚡ เปลี่ยน Status", [
        "กด dropdown ที่ช่อง Status",
        "เลือก status ใหม่",
        "กรอกข้อมูลเพิ่ม (ถ้ามี popup)",
    ]),
    ("📌 ข้อมูลในตาราง", [
        "SLA นับวันตั้งแต่ Open",
        "PIC = TA คนรับเคส (auto-assign)",
        "คลิกแถวเพื่อดูรายละเอียด",
    ]),
]
for i, (title, items) in enumerate(cols2):
    bx = Inches(0.35 + i * 4.32)
    add_card(slide, bx, feat_y, Inches(4.15), Inches(2.8), fill=GRAY_LIGHT)
    add_text(slide, title, bx + Inches(0.15), feat_y + Inches(0.12), Inches(3.8), Inches(0.4),
             size=14, bold=True, color=GREEN)
    bullet_box(slide, items, bx + Inches(0.15), feat_y + Inches(0.55), Inches(3.8), Inches(2.0),
               size=13, icon="→")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — My Cases
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "🗂️ My Cases", "เฉพาะเคสที่คุณรับผิดชอบอยู่")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=GRAY_LIGHT)

add_card(slide, Inches(0.35), Inches(1.7), Inches(6.0), Inches(4.5))
add_text(slide, "My Cases คืออะไร?", Inches(0.55), Inches(1.85), Inches(5.5), Inches(0.45),
         size=18, bold=True, color=GREEN)
items_left = [
    "แสดงเฉพาะเคสที่ assign มาให้คุณ\n(ชื่อคุณอยู่ใน PIC column)",
    "ใช้งานแทน All Requests ได้เลยถ้า\nต้องการดูแค่งานของตัวเอง",
    "เปลี่ยน Status / กรอกข้อมูล\nได้เหมือนกันทุกอย่าง",
    "มี SLA timer บอกว่าแต่ละเคส\nใช้เวลาไปแล้วกี่วัน",
]
bullet_box(slide, items_left, Inches(0.55), Inches(2.4), Inches(5.3), Inches(3.5),
           size=14, icon="✓")

add_card(slide, Inches(6.7), Inches(1.7), Inches(6.0), Inches(4.5), fill=RGBColor(0xEC,0xFD,0xF5))
add_text(slide, "Auto-Assign คืออะไร?", Inches(6.9), Inches(1.85), Inches(5.5), Inches(0.45),
         size=18, bold=True, color=TEAL)
items_right = [
    "เมื่อ TA เปลี่ยนสถานะเคสจาก Open\nเป็น Recruiting (หรือสถานะอื่น) ครั้งแรก",
    "ระบบจะ assign เคสนั้นให้\nTA คนที่กดเปลี่ยนสถานะอัตโนมัติ",
    "ชื่อ TA จะปรากฏใน PIC column\nใน Sheets และ Dashboard",
    "ถ้า Admin อยากโอนเคส → ทำได้\nจากหน้า All Requests (Admin only)",
]
bullet_box(slide, items_right, Inches(6.9), Inches(2.4), Inches(5.5), Inches(3.5),
           size=14, color=RGBColor(0x06,0x5F,0x46), icon="✓")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Offering Flow
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "📋 Offering — วิธีบันทึกข้อมูล Offer", "เมื่อส่ง Offer ให้ผู้สมัครแล้ว")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=WHITE)

# Steps
steps = [
    ("1", "เปลี่ยน Status → Offering",
     "กดที่ dropdown Status\nในตาราง แล้วเลือก Offering", INDIGO),
    ("2", "กรอก Candidate (optional)",
     "ใส่ชื่อ-นามสกุลผู้สมัคร\nที่ได้รับ Offer", AMBER),
    ("3", "วางลิ้ง CV (optional)",
     "Copy link จาก Google Drive\n(ต้องตั้ง share → Anyone with link)", GREEN),
    ("4", "กด ยืนยัน Offering",
     "ระบบบันทึก + sync ไป\nGoogle Sheets อัตโนมัติ", TEAL),
]
for i, (num, title, desc, col) in enumerate(steps):
    bx = Inches(0.35 + i * 3.25)
    # circle number
    nc = slide.shapes.add_shape(9, bx + Inches(0.9), Inches(1.8), Inches(0.65), Inches(0.65))
    nc.fill.solid(); nc.fill.fore_color.rgb = col; nc.line.fill.background()
    add_text(slide, num, bx + Inches(0.9), Inches(1.82), Inches(0.65), Inches(0.6),
             size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, bx, Inches(2.6), Inches(3.1), Inches(0.55),
             size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_card(slide, bx + Inches(0.15), Inches(3.2), Inches(2.8), Inches(1.2), fill=GRAY_LIGHT)
    add_text(slide, desc, bx + Inches(0.25), Inches(3.3), Inches(2.6), Inches(1.0),
             size=13, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(slide, "→", bx + Inches(3.0), Inches(2.6), Inches(0.4), Inches(0.55),
                 size=20, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)

# Result box
add_card(slide, Inches(0.35), Inches(4.6), Inches(12.6), Inches(1.55), fill=RGBColor(0xEC,0xFD,0xF5))
add_text(slide, "✅  ผลลัพธ์ใน Google Sheets",
         Inches(0.55), Inches(4.72), Inches(5), Inches(0.4),
         size=15, bold=True, color=GREEN)
add_text(slide,
         'Column K (Offered Candidate) จะแสดงชื่อผู้สมัครเป็น Hyperlink สีฟ้า — '
         'คลิกได้เลยเพื่อเปิด CV ใน Google Drive',
         Inches(0.55), Inches(5.2), Inches(12), Inches(0.7),
         size=14, color=RGBColor(0x06,0x5F,0x46))

add_card(slide, Inches(0.35), Inches(6.25), Inches(12.6), Inches(0.75),
         fill=RGBColor(0xFF,0xF7,0xED))
add_text(slide,
         "⚠️  ไฟล์ CV ใน Google Drive ต้องตั้ง Share → Anyone with the link (Viewer) "
         "ก่อน ไม่งั้นคนอื่นเปิดไม่ได้",
         Inches(0.55), Inches(6.35), Inches(12), Inches(0.55),
         size=13, color=ORANGE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Onboarding
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "🟦 Onboarding — บันทึกผู้สมัครที่ตอบรับ", "ใช้เมื่อผู้สมัคร Accept Offer แล้ว")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=WHITE)

# Left: steps
add_card(slide, Inches(0.35), Inches(1.7), Inches(6.5), Inches(5.1))
add_text(slide, "ขั้นตอน", Inches(0.55), Inches(1.85), Inches(6), Inches(0.4),
         size=18, bold=True, color=TEAL)
steps2 = [
    ("เปลี่ยน Status → Onboarding",   "เลือกจาก dropdown เหมือนเดิม"),
    ("กรอกชื่อ Candidate *",           "บังคับกรอก — ต้องมีชื่อก่อนกด ยืนยัน"),
    ("เลือกวันเริ่มงาน *",              "บังคับกรอก — เลือกจาก date picker"),
    ("กด ยืนยัน Onboarding",          "ระบบ sync ไป Sheets → col P (Onboard Date)"),
]
for i, (t, d) in enumerate(steps2):
    sy = Inches(2.4 + i * 0.95)
    nc = slide.shapes.add_shape(9, Inches(0.55), sy + Inches(0.05),
                                Inches(0.42), Inches(0.42))
    nc.fill.solid(); nc.fill.fore_color.rgb = TEAL; nc.line.fill.background()
    add_text(slide, str(i+1), Inches(0.55), sy + Inches(0.06), Inches(0.42), Inches(0.38),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, t, Inches(1.1), sy, Inches(5.4), Inches(0.38),
             size=14, bold=True, color=DARK)
    add_text(slide, d, Inches(1.1), sy + Inches(0.38), Inches(5.4), Inches(0.4),
             size=12, color=GRAY)

# Right: what happens
add_card(slide, Inches(7.1), Inches(1.7), Inches(5.8), Inches(5.1), fill=RGBColor(0xF0,0xFD,0xFA))
add_text(slide, "สิ่งที่จะเกิดขึ้นอัตโนมัติ", Inches(7.3), Inches(1.85), Inches(5.4), Inches(0.4),
         size=18, bold=True, color=TEAL)
effects = [
    ("Google Sheets", "บันทึกชื่อและวันเริ่มงาน\nใน column K และ P"),
    ("Slack #internal-hc-updates", "แจ้งทีมว่าเคสนี้\nอยู่ระหว่าง Onboarding"),
    ("Audit Log", "บันทึกประวัติการเปลี่ยนสถานะ\nพร้อม timestamp"),
]
for i, (label, desc) in enumerate(effects):
    ey = Inches(2.4 + i * 1.3)
    add_rect(slide, Inches(7.3), ey, Inches(0.06), Inches(0.9), fill=TEAL)
    add_text(slide, label, Inches(7.55), ey, Inches(5), Inches(0.38),
             size=14, bold=True, color=TEAL)
    add_text(slide, desc, Inches(7.55), ey + Inches(0.38), Inches(5), Inches(0.5),
             size=12, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — JD Files
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "📁 JD Files", "ไฟล์ Job Description ของแต่ละตำแหน่ง")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=GRAY_LIGHT)

add_card(slide, Inches(0.35), Inches(1.7), Inches(6.1), Inches(5.1))
add_text(slide, "JD Files คืออะไร?", Inches(0.55), Inches(1.85), Inches(5.6), Inches(0.4),
         size=18, bold=True, color=AMBER)
items_jd = [
    "เก็บไฟล์ JD (PDF/Word) แนบกับ request แต่ละอัน",
    "Manager สามารถอัพโหลด JD ตอนยื่น request",
    "TA เปิดดูได้จากหน้านี้โดยตรง",
    "ค้นหาด้วยชื่อตำแหน่ง / แผนก ได้",
    "กด Download เพื่อโหลดไฟล์มาที่เครื่อง",
]
bullet_box(slide, items_jd, Inches(0.55), Inches(2.4), Inches(5.4), Inches(3.5),
           size=14, icon="📄")

add_card(slide, Inches(6.8), Inches(1.7), Inches(6.1), Inches(5.1), fill=RGBColor(0xFF,0xFB,0xEB))
add_text(slide, "วิธีใช้", Inches(7.0), Inches(1.85), Inches(5.6), Inches(0.4),
         size=18, bold=True, color=AMBER)
steps_jd = [
    ("ไปที่เมนู JD Files", "จาก Sidebar ด้านซ้าย"),
    ("ค้นหาตำแหน่งที่ต้องการ", "พิมพ์ชื่อในช่อง Search"),
    ("คลิก Preview หรือ Download", "เปิดดู JD ออนไลน์ หรือโหลดไฟล์"),
]
for i, (t, d) in enumerate(steps_jd):
    sy = Inches(2.4 + i * 1.2)
    nc = slide.shapes.add_shape(9, Inches(7.0), sy + Inches(0.06),
                                Inches(0.42), Inches(0.42))
    nc.fill.solid(); nc.fill.fore_color.rgb = AMBER; nc.line.fill.background()
    add_text(slide, str(i+1), Inches(7.0), sy + Inches(0.07), Inches(0.42), Inches(0.38),
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, t, Inches(7.55), sy, Inches(5.1), Inches(0.38),
             size=14, bold=True, color=DARK)
    add_text(slide, d, Inches(7.55), sy + Inches(0.38), Inches(5.1), Inches(0.4),
             size=12, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Audit Log
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "🔍 Audit Log", "ประวัติทุกการเปลี่ยนแปลงในระบบ")
add_rect(slide, 0, Inches(1.45), W, H - Inches(1.45), fill=WHITE)

add_card(slide, Inches(0.35), Inches(1.7), Inches(6.1), Inches(4.8))
add_text(slide, "Audit Log บันทึกอะไรบ้าง?", Inches(0.55), Inches(1.85), Inches(5.6), Inches(0.4),
         size=18, bold=True, color=ORANGE)
items_audit = [
    "ทุกครั้งที่ Status เปลี่ยน (ใคร, เมื่อไหร่, จากอะไร → อะไร)",
    "การ Submit request ใหม่ (Manager)",
    "การ Assign TA ให้เคส",
    "การ Reject พร้อมเหตุผล",
    "การ Upload / ลบไฟล์ JD",
]
bullet_box(slide, items_audit, Inches(0.55), Inches(2.4), Inches(5.4), Inches(3.5),
           size=14, icon="📌")

add_card(slide, Inches(6.8), Inches(1.7), Inches(6.1), Inches(4.8), fill=RGBColor(0xFF,0xF7,0xED))
add_text(slide, "ใช้ประโยชน์อย่างไร?", Inches(7.0), Inches(1.85), Inches(5.6), Inches(0.4),
         size=18, bold=True, color=ORANGE)
items_use = [
    "ตรวจสอบว่าใคร Update เคสล่าสุด",
    "ดูว่าเคสนี้อยู่ใน Status นี้มานานแค่ไหน",
    "ใช้ตรวจสอบเมื่อมีข้อโต้แย้ง",
    "Export ข้อมูลสำหรับ HR Report",
    "กรองด้วย HCID / ชื่อ TA / วันที่",
]
bullet_box(slide, items_use, Inches(7.0), Inches(2.4), Inches(5.6), Inches(3.5),
           size=14, color=RGBColor(0x78,0x35,0x0F), icon="→")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Tips & Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill=GREEN)

add_text(slide, "สรุป Tips สำหรับ TA", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "จำ 5 ข้อนี้ไว้ก็พอ 😊",
         Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         size=18, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER, italic=True)

tips_final = [
    ("🔁", "เปลี่ยน Status ทันที",
     "พอมีความคืบหน้าแต่ละขั้น อย่าลืมอัพ status\nในระบบ — ทีมทั้งหมดจะเห็นได้เลย"),
    ("📋", "Offering → กรอก CV Link",
     "วาง Google Drive link ของ CV\n(ต้องตั้ง 'Anyone with link' ก่อน)"),
    ("🟦", "Onboarding → ใส่วันเริ่มงาน",
     "กรอกวันที่ถูกต้อง — จะ sync ไป Sheets\nโดยอัตโนมัติ"),
    ("🔍", "My Cases = หน้าทำงานหลัก",
     "ใช้ My Cases แทน All Requests\nเพื่อดูแค่เคสของตัวเอง"),
    ("📣", "Slack แจ้งทุก Update",
     "ทุก status change จะแจ้ง\n#internal-hc-updates อัตโนมัติ"),
]
for i, (icon, title, desc) in enumerate(tips_final):
    col_i = i % 3
    row_i = i // 3
    bx = Inches(0.35 + col_i * 4.35)
    by = Inches(1.85 + row_i * 2.3)
    card = add_rect(slide, bx, by, Inches(4.15), Inches(2.05), fill=RGBColor(0x06,0x6E,0x56))
    add_text(slide, icon + "  " + title, bx + Inches(0.18), by + Inches(0.15),
             Inches(3.8), Inches(0.5), size=16, bold=True, color=WHITE)
    add_text(slide, desc, bx + Inches(0.18), by + Inches(0.7),
             Inches(3.8), Inches(1.2), size=13, color=RGBColor(0xBB,0xF7,0xD0))

add_text(slide, "hcrequest.web.app  |  ติดต่อ Admin ถ้ามีปัญหา",
         Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
         size=13, color=RGBColor(0x6E,0xE7,0xB7), align=PP_ALIGN.CENTER, italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/flook/work/Projects/hc-request-app/HC_Request_TA_Guide.pptx"
prs.save(out)
print("Saved:", out)
