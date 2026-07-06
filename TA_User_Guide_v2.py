"""
สร้างสไลด์คู่มือ TA — ใช้รูปจริงจาก Web App
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Colors ────────────────────────────────────────────────────────────────────
GREEN        = RGBColor(0x00, 0x80, 0x65)
GREEN_DARK   = RGBColor(0x00, 0x5C, 0x49)
GREEN_LIGHT  = RGBColor(0xEC, 0xFD, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1E, 0x29, 0x3B)
GRAY         = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
INDIGO       = RGBColor(0x43, 0x38, 0xCA)
TEAL         = RGBColor(0x0D, 0x94, 0x88)
ORANGE       = RGBColor(0xEA, 0x58, 0x0C)
AMBER        = RGBColor(0xD9, 0x77, 0x06)
RED          = RGBColor(0xDC, 0x26, 0x26)
YELLOW_LIGHT = RGBColor(0xFF, 0xF9, 0xC4)
YELLOW_DARK  = RGBColor(0x92, 0x70, 0x0A)

W = Inches(13.33)
H = Inches(7.5)

SDIR = '/Users/flook/work/Projects/hc-request-app/slides_screenshots'

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill=GREEN, line_color=None, line_w=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w or 1)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size=16, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    return tb

def bullets(slide, items, x, y, w, h, size=14, color=DARK, icon='•', spacing=5):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(spacing)
        r = p.add_run(); r.text = f'{icon}  {item}'
        r.font.size = Pt(size); r.font.color.rgb = color

def header(slide, title, subtitle=None, color=GREEN):
    rect(slide, 0, 0, W, Inches(1.3), fill=color)
    txt(slide, title, Inches(0.45), Inches(0.1), Inches(10), Inches(0.75),
        size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, Inches(0.45), Inches(0.82), Inches(10), Inches(0.38),
            size=14, color=RGBColor(0xBB,0xF7,0xD0), italic=True)

def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, w, h)

def callout_box(slide, text, x, y, w=Inches(2.8), color=AMBER, bg=YELLOW_LIGHT):
    """กล่อง callout สีเหลือง"""
    r = rect(slide, x, y, w, Inches(0.48), fill=bg,
             line_color=color, line_w=1.5)
    txt(slide, text, x+Inches(0.1), y+Inches(0.05), w-Inches(0.2), Inches(0.38),
        size=12, color=YELLOW_DARK, bold=True)
    return r

def screenshot_frame(slide, img_path, x, y, w, h):
    """ใส่รูปพร้อม border shadow"""
    # shadow
    rect(slide, x+Inches(0.06), y+Inches(0.06), w, h,
         fill=RGBColor(0xCB,0xD5,0xE1))
    # white border
    r = rect(slide, x, y, w, h, fill=WHITE,
             line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
    add_image(slide, img_path, x+Inches(0.03), y+Inches(0.03),
              w-Inches(0.06), h-Inches(0.06))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GREEN)
rect(slide, 0, 0, W, Inches(0.06), fill=GREEN_DARK)

# Dashboard screenshot as bg (faded right side)
img_path = f'{SDIR}/crop_01_dashboard.png'
if os.path.exists(img_path):
    pic = add_image(slide, img_path, Inches(5.5), Inches(0.8), Inches(7.5), Inches(6.5))
    # overlay green fade on left
    rect(slide, 0, 0, Inches(7), H, fill=GREEN)
    # semi-transparent overlay on image area
    ov = rect(slide, Inches(5.5), Inches(0.8), Inches(7.5), Inches(6.5),
              fill=RGBColor(0x00,0x70,0x58))
    sp = ov.element
    sp_pr = sp.find(qn('p:spPr'))
    # set alpha on solid fill
    solid = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
    if solid is not None:
        srgb = solid.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            alpha_el = etree.SubElement(srgb,
                '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
            alpha_el.set('val', '75000')  # 75% opacity

txt(slide, '📋  HC Request System', Inches(0.5), Inches(0.9), Inches(6.5), Inches(0.7),
    size=20, color=RGBColor(0xBB,0xF7,0xD0), italic=True)
txt(slide, 'คู่มือสำหรับ\nTA Team', Inches(0.5), Inches(1.6), Inches(6.5), Inches(2.4),
    size=52, bold=True, color=WHITE)
txt(slide, 'วิธีรับเคส • เปลี่ยน Status • จัดการข้อมูล Candidate',
    Inches(0.5), Inches(4.15), Inches(6.5), Inches(0.6),
    size=17, color=RGBColor(0xBB,0xF7,0xD0))

rect(slide, Inches(0.5), Inches(4.9), Inches(2.4), Inches(0.04),
     fill=RGBColor(0x6E,0xE7,0xB7))

txt(slide, 'People Experience · Freshket', Inches(0.5), Inches(5.1),
    Inches(6), Inches(0.4), size=14, color=RGBColor(0x6E,0xE7,0xB7))
txt(slide, 'hcrequest.web.app', Inches(0.5), Inches(5.5),
    Inches(6), Inches(0.35), size=13, color=RGBColor(0x6E,0xE7,0xB7), italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Dashboard
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
header(slide, '📊  Dashboard', 'หน้าแรกหลัง Login — ภาพรวมทุกเคส')

# Screenshot ใหญ่ซ้าย
screenshot_frame(slide, f'{SDIR}/crop_01_dashboard.png',
                 Inches(0.3), Inches(1.4), Inches(8.5), Inches(5.8))

# Panel ขวา
rect(slide, Inches(9.1), Inches(1.4), Inches(4.0), Inches(5.8), fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)

txt(slide, 'สิ่งที่เห็นใน Dashboard', Inches(9.25), Inches(1.55),
    Inches(3.7), Inches(0.45), size=15, bold=True, color=GREEN)

items = [
    ('📦 Stat Cards', 'จำนวนเคสแต่ละสถานะ (Open, In Progress, Offering…)'),
    ('📅 Monthly Chart', 'กราฟ HC Request รายเดือน เทียบปีก่อน'),
    ('📊 Breakdown', 'แยกตามแผนก / ตำแหน่ง 6 เดือนล่าสุด'),
    ('🔄 รายการ', 'เปลี่ยน Status และดู SLA ได้จากตาราง'),
]
for i, (title, desc) in enumerate(items):
    y = Inches(2.1 + i * 1.1)
    rect(slide, Inches(9.25), y, Inches(0.05), Inches(0.8), fill=GREEN)
    txt(slide, title, Inches(9.45), y, Inches(3.4), Inches(0.38),
        size=13, bold=True, color=GREEN)
    txt(slide, desc, Inches(9.45), y+Inches(0.38), Inches(3.4), Inches(0.5),
        size=12, color=GRAY)

callout_box(slide, '💡 กด "รายการ" มุมขวาบน\nเพื่อดูตาราง request ทั้งหมด',
            Inches(9.25), Inches(6.5), w=Inches(3.7))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — All Requests
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
header(slide, '📋  All Requests', 'ดู request ทั้งหมด — เปลี่ยน Status ได้จากหน้านี้')

# Screenshot
screenshot_frame(slide, f'{SDIR}/crop_02_all_requests.png',
                 Inches(0.3), Inches(1.4), Inches(8.5), Inches(5.8))

# Callouts overlay บนรูป
callout_box(slide, '① Status Tabs — กรองด้วยสถานะ',
            Inches(0.4), Inches(2.0), color=INDIGO,
            bg=RGBColor(0xEE,0xF2,0xFF))
callout_box(slide, '② Dropdown Status — คลิกเพื่อเปลี่ยน',
            Inches(0.4), Inches(3.0), color=GREEN, bg=GREEN_LIGHT)
callout_box(slide, '③ SLA — วันที่ใช้ไปแต่ละเคส',
            Inches(0.4), Inches(4.0), color=ORANGE,
            bg=RGBColor(0xFF,0xF7,0xED))
callout_box(slide, '④ ค้นหา / Filter ด้านบน',
            Inches(0.4), Inches(5.0), color=TEAL,
            bg=RGBColor(0xF0,0xFD,0xFA))

# Panel ขวา
rect(slide, Inches(9.1), Inches(1.4), Inches(4.0), Inches(5.8), fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'วิธีเปลี่ยน Status', Inches(9.25), Inches(1.55),
    Inches(3.7), Inches(0.45), size=15, bold=True, color=GREEN)

steps = [
    ('1', 'หาเคสที่ต้องการ', 'ค้นหาหรือ filter'),
    ('2', 'คลิก Dropdown', 'ที่ช่อง Status'),
    ('3', 'เลือก Status ใหม่', ''),
    ('4', 'กรอกข้อมูล (ถ้ามี)', 'Popup จะขึ้นเอง'),
    ('5', 'กด ยืนยัน', 'ระบบ sync อัตโนมัติ'),
]
for i, (n, t, d) in enumerate(steps):
    y = Inches(2.1 + i * 0.9)
    nc = slide.shapes.add_shape(9, Inches(9.25), y+Inches(0.05),
                                Inches(0.4), Inches(0.4))
    nc.fill.solid(); nc.fill.fore_color.rgb = GREEN
    nc.line.fill.background()
    txt(slide, n, Inches(9.25), y+Inches(0.06), Inches(0.4), Inches(0.38),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, t, Inches(9.78), y, Inches(3.1), Inches(0.35),
        size=13, bold=True, color=DARK)
    if d:
        txt(slide, d, Inches(9.78), y+Inches(0.34), Inches(3.1), Inches(0.3),
            size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — My Cases
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
header(slide, '🗂️  My Cases', 'เฉพาะเคสที่คุณรับผิดชอบ')

screenshot_frame(slide, f'{SDIR}/crop_03_my_cases.png',
                 Inches(4.0), Inches(1.4), Inches(9.05), Inches(5.8))

# Left panel
rect(slide, Inches(0.3), Inches(1.4), Inches(3.5), Inches(5.8), fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'My Cases\nคืออะไร?', Inches(0.45), Inches(1.55),
    Inches(3.2), Inches(0.85), size=18, bold=True, color=GREEN)

points = [
    'แสดงเฉพาะเคสที่\nระบบ assign มาให้คุณ',
    'ใช้แทน All Requests\nสำหรับงานของตัวเอง',
    'เปลี่ยน Status ได้\nเหมือนกันทุกอย่าง',
    'มี SLA Timer บอก\nวันที่ผ่านมาแต่ละเคส',
]
for i, p in enumerate(points):
    y = Inches(2.55 + i * 1.0)
    nc = slide.shapes.add_shape(9, Inches(0.45), y, Inches(0.35), Inches(0.35))
    nc.fill.solid(); nc.fill.fore_color.rgb = TEAL
    nc.line.fill.background()
    txt(slide, '✓', Inches(0.45), y+Inches(0.02), Inches(0.35), Inches(0.32),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, p, Inches(0.9), y, Inches(2.7), Inches(0.75),
        size=13, color=DARK)

rect(slide, Inches(0.45), Inches(6.55), Inches(3.2), Inches(0.5),
     fill=GREEN_LIGHT, line_color=GREEN, line_w=1)
txt(slide, '🤖 Auto-Assign: เมื่อ TA เปลี่ยน status\nจาก Open ครั้งแรก → ระบบ assign ให้คุณอัตโนมัติ',
    Inches(0.55), Inches(6.6), Inches(3.0), Inches(0.6),
    size=11, color=GREEN_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Status Workflow (visual)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=WHITE)
header(slide, '🔄  Status Workflow', 'ลำดับขั้นตอนการเปลี่ยนสถานะ')

statuses = [
    ('Open',        '📂', GRAY,   'Manager เปิดเคสใหม่\nรอ TA รับ'),
    ('Recruiting',  '🔍', GREEN,  'TA รับเคส\nกำลังหาผู้สมัคร'),
    ('Interviewing','🗣️', INDIGO, 'อยู่ระหว่าง\nสัมภาษณ์'),
    ('Offering',    '📋', AMBER,  'ส่ง Offer\nบันทึกชื่อ + CV'),
    ('Onboarding',  '🟦', TEAL,  'รอเริ่มงาน\nบันทึกวันที่'),
    ('Closed',      '✅', RGBColor(0x16,0xA3,0x4A), 'Onboard\nสำเร็จ! 🎉'),
]

bw = Inches(1.82)
bh = Inches(2.3)
gap = Inches(0.28)
sx = Inches(0.28)
sy = Inches(1.75)

for i, (name, icon, col, desc) in enumerate(statuses):
    bx = sx + i*(bw+gap)
    # shadow
    rect(slide, bx+Inches(0.05), sy+Inches(0.05), bw, bh,
         fill=RGBColor(0xE2,0xE8,0xF0))
    # card
    card = rect(slide, bx, sy, bw, bh, fill=WHITE,
                line_color=col, line_w=2)
    rect(slide, bx, sy, bw, Inches(0.1), fill=col)
    txt(slide, icon, bx, sy+Inches(0.18), bw, Inches(0.65),
        size=30, align=PP_ALIGN.CENTER)
    txt(slide, name, bx, sy+Inches(0.9), bw, Inches(0.45),
        size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
    txt(slide, desc, bx+Inches(0.1), sy+Inches(1.38), bw-Inches(0.2), Inches(0.8),
        size=12, color=GRAY, align=PP_ALIGN.CENTER)
    if i < len(statuses)-1:
        ax = bx+bw+Inches(0.04)
        txt(slide, '→', ax, sy+Inches(0.8), gap, Inches(0.55),
            size=20, bold=True, color=RGBColor(0xCB,0xD5,0xE1),
            align=PP_ALIGN.CENTER)

# Note boxes
rect(slide, Inches(0.28), Inches(4.3), Inches(6.1), Inches(0.75),
     fill=RGBColor(0xFE,0xF2,0xF2),
     line_color=RGBColor(0xFC,0xA5,0xA5), line_w=1)
txt(slide, '❌  Rejected — TA สามารถ Reject ได้ทุกสถานะ พร้อมกรอกเหตุผล\n'
    '    ถ้ายังหาคนไม่ได้ → กลับมา Recruit ใหม่ได้เสมอ',
    Inches(0.45), Inches(4.38), Inches(5.8), Inches(0.6), size=13, color=RED)

rect(slide, Inches(6.8), Inches(4.3), Inches(6.2), Inches(0.75),
     fill=RGBColor(0xFF,0xF7,0xED),
     line_color=RGBColor(0xFB,0xBF,0x24), line_w=1)
txt(slide, '📋 Offering → กรอกชื่อ Candidate + ลิ้ง CV (Google Drive)\n'
    '🟦 Onboarding → กรอก ชื่อ Candidate + วันเริ่มงาน (บังคับ)',
    Inches(6.95), Inches(4.38), Inches(5.9), Inches(0.6), size=13, color=AMBER)

rect(slide, Inches(0.28), Inches(5.2), Inches(12.75), Inches(1.5),
     fill=RGBColor(0xF8,0xFA,0xFF),
     line_color=RGBColor(0xC7,0xD2,0xFE), line_w=1)
txt(slide, '🔔  Slack แจ้งเตือนอัตโนมัติ', Inches(0.45), Inches(5.32),
    Inches(6), Inches(0.38), size=14, bold=True, color=INDIGO)
txt(slide,
    '#internal-hc-request  →  เมื่อ Manager เปิดเคสใหม่\n'
    '#internal-hc-updates   →  ทุกครั้งที่ Status เปลี่ยน',
    Inches(0.45), Inches(5.72), Inches(12.3), Inches(0.75),
    size=13, color=RGBColor(0x44,0x38,0xCA))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Offering (text + diagram เพราะ modal ถ่ายรูปไม่ได้)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=WHITE)
header(slide, '📋  Offering — บันทึก Candidate + CV', 'เมื่อส่ง Offer ให้ผู้สมัครแล้ว', color=AMBER)

# Modal mockup
modal_x = Inches(0.8)
modal_y = Inches(1.5)
modal_w = Inches(5.2)
modal_h = Inches(5.5)

rect(slide, modal_x, modal_y, modal_w, modal_h, fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=2)
rect(slide, modal_x, modal_y, modal_w, Inches(0.06),
     fill=AMBER)

txt(slide, '📋  Offering', modal_x+Inches(0.3), modal_y+Inches(0.2),
    Inches(4.5), Inches(0.5), size=22, bold=True, color=DARK)
txt(slide, 'กรอกชื่อผู้สมัครที่ได้รับ offer',
    modal_x+Inches(0.3), modal_y+Inches(0.7), Inches(4.5), Inches(0.35),
    size=14, color=GRAY)

# Candidate field
txt(slide, 'ชื่อ CANDIDATE (OPTIONAL)',
    modal_x+Inches(0.3), modal_y+Inches(1.2), Inches(4.5), Inches(0.3),
    size=10, bold=True, color=GRAY)
rect(slide, modal_x+Inches(0.3), modal_y+Inches(1.55), Inches(4.55), Inches(0.5),
     fill=WHITE, line_color=RGBColor(0xA5,0xB4,0xFC), line_w=1.5)
txt(slide, 'ชื่อ-นามสกุล ผู้สมัคร',
    modal_x+Inches(0.45), modal_y+Inches(1.62), Inches(4), Inches(0.35),
    size=13, color=RGBColor(0xCB,0xD5,0xE1))

# CV URL field
txt(slide, 'ลิ้ง CV (OPTIONAL)',
    modal_x+Inches(0.3), modal_y+Inches(2.3), Inches(4.5), Inches(0.3),
    size=10, bold=True, color=GRAY)
rect(slide, modal_x+Inches(0.3), modal_y+Inches(2.65), Inches(4.55), Inches(0.5),
     fill=WHITE, line_color=RGBColor(0xA5,0xB4,0xFC), line_w=1.5)
txt(slide, 'https://drive.google.com/...',
    modal_x+Inches(0.45), modal_y+Inches(2.72), Inches(4), Inches(0.35),
    size=13, color=RGBColor(0xCB,0xD5,0xE1))

# Buttons
rect(slide, modal_x+Inches(0.3), modal_y+Inches(4.4), Inches(1.9), Inches(0.55),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'ยกเลิก',
    modal_x+Inches(0.3), modal_y+Inches(4.47), Inches(1.9), Inches(0.42),
    size=13, color=GRAY, align=PP_ALIGN.CENTER)
rect(slide, modal_x+Inches(2.6), modal_y+Inches(4.4), Inches(2.3), Inches(0.55),
     fill=INDIGO)
txt(slide, 'ยืนยัน Offering',
    modal_x+Inches(2.6), modal_y+Inches(4.47), Inches(2.3), Inches(0.42),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right: explanation
rx = Inches(6.5)
txt(slide, 'วิธีใช้ Offering Modal',
    rx, Inches(1.5), Inches(6.5), Inches(0.5),
    size=18, bold=True, color=AMBER)

steps = [
    ('ชื่อ Candidate', 'ใส่ชื่อ-นามสกุลผู้สมัคร\n(กรอกหรือไม่ก็ได้)'),
    ('ลิ้ง CV', 'วาง Google Drive link ของ CV\n⚠️ ต้องตั้ง Share "Anyone with link" ก่อน'),
    ('ยืนยัน Offering', 'กดปุ่ม → ระบบบันทึก + sync\nไป Google Sheets อัตโนมัติ'),
]
for i, (t, d) in enumerate(steps):
    sy2 = Inches(2.2 + i*1.35)
    nc = slide.shapes.add_shape(9, rx, sy2+Inches(0.04),
                                Inches(0.42), Inches(0.42))
    nc.fill.solid(); nc.fill.fore_color.rgb = AMBER
    nc.line.fill.background()
    txt(slide, str(i+1), rx, sy2+Inches(0.05), Inches(0.42), Inches(0.38),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, t, rx+Inches(0.58), sy2, Inches(5.7), Inches(0.38),
        size=14, bold=True, color=DARK)
    txt(slide, d, rx+Inches(0.58), sy2+Inches(0.38), Inches(5.7), Inches(0.65),
        size=13, color=GRAY)

rect(slide, rx, Inches(6.25), Inches(6.55), Inches(0.7),
     fill=GREEN_LIGHT, line_color=GREEN, line_w=1)
txt(slide, '✅ ผลลัพธ์ใน Sheets: Column "Offered Candidate" จะกลายเป็น\n'
    'Hyperlink สีฟ้า คลิกเปิด CV ใน Google Drive ได้เลย',
    rx+Inches(0.15), Inches(6.32), Inches(6.2), Inches(0.6),
    size=12, color=GREEN_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Onboarding (modal mockup)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=WHITE)
header(slide, '🟦  Onboarding — บันทึกวันเริ่มงาน', 'เมื่อผู้สมัคร Accept Offer', color=TEAL)

# Modal mockup
modal_x = Inches(0.8)
modal_y = Inches(1.5)
modal_w = Inches(5.2)
modal_h = Inches(5.5)

rect(slide, modal_x, modal_y, modal_w, modal_h, fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=2)
rect(slide, modal_x, modal_y, modal_w, Inches(0.06), fill=TEAL)

txt(slide, '🟦  Waiting Onboarding', modal_x+Inches(0.3), modal_y+Inches(0.2),
    Inches(4.5), Inches(0.5), size=20, bold=True, color=DARK)
txt(slide, 'กรุณากรอกข้อมูลผู้สมัครที่รับ offer',
    modal_x+Inches(0.3), modal_y+Inches(0.7), Inches(4.5), Inches(0.35),
    size=14, color=GRAY)

# Candidate (required)
txt(slide, 'ชื่อ CANDIDATE *',
    modal_x+Inches(0.3), modal_y+Inches(1.2), Inches(4.5), Inches(0.3),
    size=10, bold=True, color=RED)
rect(slide, modal_x+Inches(0.3), modal_y+Inches(1.55), Inches(4.55), Inches(0.5),
     fill=WHITE, line_color=RGBColor(0x5E,0xEA,0xD4), line_w=1.5)
txt(slide, 'ชื่อ-นามสกุล ผู้สมัคร *',
    modal_x+Inches(0.45), modal_y+Inches(1.62), Inches(4), Inches(0.35),
    size=13, color=RGBColor(0xCB,0xD5,0xE1))

# Start date (required)
txt(slide, 'วันเริ่มงาน *',
    modal_x+Inches(0.3), modal_y+Inches(2.3), Inches(4.5), Inches(0.3),
    size=10, bold=True, color=RED)
rect(slide, modal_x+Inches(0.3), modal_y+Inches(2.65), Inches(4.55), Inches(0.5),
     fill=WHITE, line_color=RGBColor(0x5E,0xEA,0xD4), line_w=1.5)
txt(slide, '📅  เลือกวันที่',
    modal_x+Inches(0.45), modal_y+Inches(2.72), Inches(4), Inches(0.35),
    size=13, color=RGBColor(0xCB,0xD5,0xE1))

rect(slide, modal_x+Inches(0.3), modal_y+Inches(3.35), Inches(4.55), Inches(0.45),
     fill=RGBColor(0xFF,0xFB,0xEB), line_color=AMBER, line_w=1)
txt(slide, '⚠️ ทั้งสองช่องนี้ บังคับกรอก ก่อนกด ยืนยัน',
    modal_x+Inches(0.45), modal_y+Inches(3.42), Inches(4.1), Inches(0.32),
    size=12, color=AMBER)

# Button
rect(slide, modal_x+Inches(0.3), modal_y+Inches(4.4), Inches(1.9), Inches(0.55),
     fill=WHITE, line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'ยกเลิก', modal_x+Inches(0.3), modal_y+Inches(4.47),
    Inches(1.9), Inches(0.42), size=13, color=GRAY, align=PP_ALIGN.CENTER)
rect(slide, modal_x+Inches(2.6), modal_y+Inches(4.4), Inches(2.3), Inches(0.55),
     fill=TEAL)
txt(slide, 'ยืนยัน Onboarding',
    modal_x+Inches(2.6), modal_y+Inches(4.47), Inches(2.3), Inches(0.42),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Right explanation
rx = Inches(6.5)
txt(slide, 'สิ่งที่เกิดขึ้นอัตโนมัติ',
    rx, Inches(1.5), Inches(6.5), Inches(0.5),
    size=18, bold=True, color=TEAL)

effects = [
    ('📊 Google Sheets', 'บันทึกชื่อ Candidate ใน Column K\nและวันเริ่มงานใน Column P'),
    ('📣 Slack แจ้งทีม', 'ส่งข้อความไป #internal-hc-updates\nว่าเคสนี้อยู่ระหว่าง Onboarding'),
    ('📜 Audit Log', 'บันทึกประวัติการเปลี่ยน Status\nพร้อม timestamp อัตโนมัติ'),
]
for i, (label, desc) in enumerate(effects):
    ey = Inches(2.2 + i*1.4)
    rect(slide, rx, ey, Inches(0.06), Inches(1.0), fill=TEAL)
    txt(slide, label, rx+Inches(0.22), ey, Inches(6.0), Inches(0.42),
        size=14, bold=True, color=TEAL)
    txt(slide, desc, rx+Inches(0.22), ey+Inches(0.42), Inches(6.0), Inches(0.65),
        size=13, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — JD Files
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
header(slide, '📁  JD Files', 'ไฟล์ Job Description ของแต่ละตำแหน่ง', color=AMBER)

screenshot_frame(slide, f'{SDIR}/crop_04_jd_files.png',
                 Inches(0.3), Inches(1.4), Inches(8.5), Inches(5.8))

rect(slide, Inches(9.1), Inches(1.4), Inches(4.0), Inches(5.8), fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'JD Files ใช้ทำอะไร?', Inches(9.25), Inches(1.55),
    Inches(3.7), Inches(0.45), size=15, bold=True, color=AMBER)

items_jd = [
    ('📄 ดู JD', 'เปิดอ่านรายละเอียดงาน\nก่อนเริ่ม Recruit'),
    ('🔍 ค้นหา', 'หาด้วยชื่อตำแหน่ง\nหรือชื่อแผนก'),
    ('⬇️ Download', 'โหลดไฟล์มาที่เครื่อง\nหรือ Share ต่อได้'),
    ('📋 Preview', 'ดูออนไลน์โดยตรง\nไม่ต้อง download'),
]
for i, (t, d) in enumerate(items_jd):
    y = Inches(2.15 + i * 1.07)
    rect(slide, Inches(9.25), y, Inches(0.05), Inches(0.8), fill=AMBER)
    txt(slide, t, Inches(9.45), y, Inches(3.4), Inches(0.38),
        size=13, bold=True, color=AMBER)
    txt(slide, d, Inches(9.45), y+Inches(0.38), Inches(3.4), Inches(0.5),
        size=12, color=GRAY)

callout_box(slide, '💡 Manager จะ upload JD\nตอนยื่น Request เอง',
            Inches(9.25), Inches(6.5), w=Inches(3.7), color=AMBER,
            bg=RGBColor(0xFF,0xFB,0xEB))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Audit Log
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GRAY_LIGHT)
header(slide, '🔍  Audit Log', 'ประวัติทุกการเปลี่ยนแปลงในระบบ', color=ORANGE)

screenshot_frame(slide, f'{SDIR}/crop_05_audit_log.png',
                 Inches(4.0), Inches(1.4), Inches(9.05), Inches(5.8))

rect(slide, Inches(0.3), Inches(1.4), Inches(3.5), Inches(5.8), fill=WHITE,
     line_color=RGBColor(0xE2,0xE8,0xF0), line_w=1)
txt(slide, 'บันทึกอะไร\nบ้าง?', Inches(0.45), Inches(1.55),
    Inches(3.2), Inches(0.85), size=18, bold=True, color=ORANGE)

logged = [
    'ทุกครั้งที่ Status เปลี่ยน',
    'Manager Submit request ใหม่',
    'TA รับ / โอนเคส',
    'Reject พร้อมเหตุผล',
    'Upload / ลบ JD',
]
for i, p in enumerate(logged):
    y = Inches(2.6 + i * 0.72)
    rect(slide, Inches(0.45), y+Inches(0.1), Inches(0.06), Inches(0.35), fill=ORANGE)
    txt(slide, p, Inches(0.65), y, Inches(2.8), Inches(0.5),
        size=13, color=DARK)

rect(slide, Inches(0.45), Inches(6.45), Inches(3.2), Inches(0.55),
     fill=RGBColor(0xFF,0xF7,0xED), line_color=ORANGE, line_w=1)
txt(slide, '📊 กรอง / Export ได้\nตาม HCID, TA, หรือวันที่',
    Inches(0.55), Inches(6.5), Inches(3.0), Inches(0.55),
    size=12, color=ORANGE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Tips Summary
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, W, H, fill=GREEN)
rect(slide, 0, 0, W, Inches(0.06), fill=GREEN_DARK)

txt(slide, 'สรุป Tips สำหรับ TA 🚀',
    Inches(0.5), Inches(0.25), Inches(12), Inches(0.75),
    size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(slide, 'จำ 5 ข้อนี้ไว้ก็พอ',
    Inches(0.5), Inches(1.0), Inches(12), Inches(0.45),
    size=17, color=RGBColor(0xBB,0xF7,0xD0), align=PP_ALIGN.CENTER, italic=True)

tips = [
    ('🔁', 'อัพ Status ทันที',
     'พอมีความคืบหน้า อย่ารอ —\nทีมทั้งหมดจะเห็นได้เลย'),
    ('📋', 'Offering → ใส่ CV Link',
     'วาง Google Drive link\n(ต้อง Share "Anyone with link")'),
    ('🟦', 'Onboarding → วันเริ่มงาน',
     'กรอกวันที่ให้ถูกต้อง\nSync ไป Sheets อัตโนมัติ'),
    ('🗂️', 'My Cases = หน้าทำงาน',
     'ดูแค่เคสของตัวเอง\nสะดวกกว่า All Requests'),
    ('📣', 'Slack แจ้งทุก Update',
     'ทุก Status change แจ้ง\n#internal-hc-updates เอง'),
]
for i, (icon, title, desc) in enumerate(tips):
    col_i = i % 3
    row_i = i // 3
    bx = Inches(0.55 + col_i * 4.22)
    by = Inches(1.65 + row_i * 2.35)
    w2 = Inches(3.95)
    h2 = Inches(2.15)
    card = rect(slide, bx, by, w2, h2, fill=GREEN_DARK)
    rect(slide, bx, by, w2, Inches(0.07), fill=RGBColor(0x6E,0xE7,0xB7))
    txt(slide, icon + '  ' + title,
        bx+Inches(0.2), by+Inches(0.2), w2-Inches(0.4), Inches(0.5),
        size=16, bold=True, color=WHITE)
    txt(slide, desc, bx+Inches(0.2), by+Inches(0.8),
        w2-Inches(0.4), Inches(1.15),
        size=14, color=RGBColor(0xBB,0xF7,0xD0))

txt(slide, 'hcrequest.web.app  |  ติดต่อ Admin หากมีปัญหา',
    Inches(0.5), Inches(7.12), Inches(12), Inches(0.35),
    size=13, color=RGBColor(0x6E,0xE7,0xB7), align=PP_ALIGN.CENTER, italic=True)

# ── Save ──────────────────────────────────────────────────────────────────────
out = '/Users/flook/work/Projects/hc-request-app/HC_Request_TA_Guide_v2.pptx'
prs.save(out)
print('Saved:', out)
